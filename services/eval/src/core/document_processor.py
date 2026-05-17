"""Office 和 PDF 文件文本提取工具。"""
import io
import logging
import base64
import tempfile
import os
from typing import Optional, List

logger = logging.getLogger(__name__)

def extract_frames_from_video(content_bytes: bytes, num_frames: int = 8) -> List[str]:
    """从视频字节流中提取 N 帧，返回 base64 列表。"""
    try:
        import cv2
        
        frames_base64 = []
        # 使用 tempfile 处理，因为 cv2.VideoCapture 需要文件路径
        with tempfile.NamedTemporaryFile(suffix=".tmp_video", delete=False) as tmp:
            tmp.write(content_bytes)
            tmp_path = tmp.name
        
        try:
            cap = cv2.VideoCapture(tmp_path)
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            if total_frames <= 0:
                logger.warning("Video has no frames or invalid format")
                return []
            
            # 均匀采样 N 帧
            if total_frames < num_frames:
                indices = list(range(total_frames))
            else:
                indices = [int(i * (total_frames - 1) / (num_frames - 1)) for i in range(num_frames)]
            
            for idx in indices:
                cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
                ret, frame = cap.read()
                if ret:
                    # 转换为 JPEG base64
                    success, buffer = cv2.imencode('.jpg', frame)
                    if success:
                        frames_base64.append(base64.b64encode(buffer).decode('utf-8'))
            cap.release()
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        
        return frames_base64
    except ImportError:
        logger.error("opencv-python-headless not installed, cannot process video")
        return []
    except Exception as e:
        logger.error(f"Failed to extract frames from video: {e}")
        return []

def extract_text_from_docx(content_bytes: bytes) -> str:
    """从 .docx 文件提取文本。"""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content_bytes))
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        logger.error(f"Failed to extract text from DOCX: {e}")
        return ""

def extract_text_from_xlsx(content_bytes: bytes) -> str:
    """从 .xlsx 文件提取文本（所有工作表）。"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content_bytes), data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            text_parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                text_parts.append(" ".join([str(cell) for cell in row if cell is not None]))
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Failed to extract text from XLSX: {e}")
        return ""

def extract_text_from_pptx(content_bytes: bytes) -> str:
    """从 .pptx 文件提取文本。"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content_bytes))
        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Failed to extract text from PPTX: {e}")
        return ""

def extract_text_from_pdf(content_bytes: bytes) -> str:
    """从 .pdf 文件提取文本。"""
    try:
        from pypdf import PdfReader
        reader = pypdf.PdfReader(io.BytesIO(content_bytes))
        text_parts = []
        for page in reader.pages:
            text_parts.append(page.extract_text())
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Failed to extract text from PDF: {e}")
        return ""

def extract_text_by_extension(extension: str, content_bytes: bytes) -> Optional[str]:
    """根据扩展名选择合适的提取方法。"""
    ext = extension.lower()
    if ext == ".docx":
        return extract_text_from_docx(content_bytes)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(content_bytes)
    elif ext == ".pptx":
        return extract_text_from_pptx(content_bytes)
    elif ext in (".pdf"):
        return extract_text_from_pdf(content_bytes)
    return None

def prepare_text_content(content_base64: str, file_path: str) -> Optional[str]:
    """
    统一的文本内容准备接口。
    封装了 base64 解码、Office/PDF 提取、纯文本回退以及过短文本的归一化处理。
    """
    if not content_base64:
        return None
        
    _, ext = os.path.splitext(file_path.lower())
    raw_bytes = base64.b64decode(content_base64)
    
    text_content = extract_text_by_extension(ext, raw_bytes)
    
    if text_content is None:
        try:
            # 尝试直接 utf-8 解码 Gitea 原生内容
            text_content = base64.b64decode(content_base64).decode("utf-8", errors="replace")
        except Exception:
            text_content = raw_bytes.decode("utf-8", errors="replace")

    if text_content and len(text_content.strip()) < 10:
        text_content = f"Text content is short: {text_content.strip()}"
        
    return text_content

def prepare_video_frames(content_base64: str, num_frames: int = 8) -> List[str]:
    """
    统一的视频帧准备接口。
    封装了 base64 解码和帧提取逻辑。
    """
    if not content_base64:
        return []
    raw_bytes = base64.b64decode(content_base64)
    return extract_frames_from_video(raw_bytes, num_frames)
